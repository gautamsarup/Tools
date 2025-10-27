#!/usr/bin/env python3
"""
PowerPoint Text Extractor with OCR, Table Extraction, and LLM Formatting

This program extracts text from PowerPoint presentations with:
- OCR support for image-based slides
- Table extraction to Excel files
- User prompts for slide selection
- LLM-powered text formatting for better readability
- Proper handling of text positioned anywhere on slides

Dependencies:
- python-pptx: For PowerPoint native text extraction
- pytesseract: For OCR processing
- Pillow: For image processing
- pandas: For table data manipulation
- openpyxl: For Excel export
- openai: For LLM text formatting
"""

import os
import sys
import argparse
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import logging
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as e:
    print(f"Missing required library: {e}")
    print("Please install required packages:")
    print("pip install python-pptx")
    sys.exit(1)

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("Note: OCR not available. Install pytesseract and Pillow for OCR support.")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    print("Note: pandas not available. Install pandas for table export.")

try:
    import openai
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False
    print("Note: OpenAI not available. Install openai for LLM formatting.")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class PowerPointExtractor:
    """
    A comprehensive PowerPoint text extractor with OCR, table extraction, and LLM formatting.
    """
    
    def __init__(self, tesseract_path: Optional[str] = None, openai_api_key: Optional[str] = None):
        """
        Initialize the PowerPoint extractor.
        
        Args:
            tesseract_path: Path to tesseract executable (auto-detected if None)
            openai_api_key: OpenAI API key for LLM formatting
        """
        self.tesseract_path = tesseract_path
        self.openai_api_key = openai_api_key
        self.openai_client = None
        self._setup_tesseract()
        self._setup_openai()
        
    def _setup_tesseract(self):
        """Setup tesseract path if provided."""
        if not OCR_AVAILABLE:
            return
            
        # Priority: 1. Constructor arg, 2. config.py file, 3. Auto-detection
        tesseract_path = self.tesseract_path
        
        if not tesseract_path:
            # Try to load from config.py
            try:
                import config
                if hasattr(config, 'TESSERACT_PATH') and config.TESSERACT_PATH:
                    tesseract_path = config.TESSERACT_PATH
                    logger.info("Using Tesseract path from config.py")
            except ImportError:
                pass
        
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            logger.info(f"Using tesseract at: {tesseract_path}")
        else:
            # Try to auto-detect tesseract
            common_paths = [
                r'C:\Program Files\Tesseract-OCR\tesseract.exe',  # Windows
                r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',  # Windows 32-bit
                '/usr/bin/tesseract',  # Linux
                '/usr/local/bin/tesseract',  # macOS
                '/opt/homebrew/bin/tesseract',  # macOS with Homebrew
            ]
            
            for path in common_paths:
                if os.path.exists(path):
                    pytesseract.pytesseract.tesseract_cmd = path
                    logger.info(f"Auto-detected tesseract at: {path}")
                    return
            
            logger.warning("Tesseract not found in common locations. Please install tesseract-ocr.")
    
    def _setup_openai(self):
        """Setup OpenAI API key."""
        if not OPENAI_AVAILABLE:
            return
        
        # Priority: 1. Constructor arg, 2. Environment variable, 3. config.py file
        api_key = None
        
        if self.openai_api_key:
            api_key = self.openai_api_key
            logger.info("OpenAI API key configured from constructor")
        elif os.getenv('OPENAI_API_KEY'):
            api_key = os.getenv('OPENAI_API_KEY')
            logger.info("Using OpenAI API key from environment variable")
        else:
            # Try to load from config.py
            try:
                import config
                if hasattr(config, 'OPENAI_API_KEY') and config.OPENAI_API_KEY:
                    api_key = config.OPENAI_API_KEY
                    logger.info("Using OpenAI API key from config.py")
            except ImportError:
                pass
        
        if api_key:
            self.openai_client = OpenAI(api_key=api_key)
            logger.info("OpenAI client initialized successfully")
        else:
            logger.warning("OpenAI API key not found. LLM formatting will be disabled.")
    
    def extract_text_from_pptx(self, pptx_path: str, slide_numbers: Optional[List[int]] = None,
                               use_ocr: bool = True, use_llm_formatting: bool = True) -> Dict[str, Any]:
        """
        Extract text from PowerPoint presentation.
        
        Args:
            pptx_path: Path to the PowerPoint file
            slide_numbers: List of slide numbers to extract (1-indexed). None for all slides.
            use_ocr: Whether to use OCR for image-based content
            use_llm_formatting: Whether to use LLM for text formatting
            
        Returns:
            Dictionary containing extracted content
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        logger.info(f"Processing PowerPoint: {pptx_path}")
        
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        
        # Validate slide numbers
        if slide_numbers:
            slide_numbers = [sn for sn in slide_numbers if 1 <= sn <= total_slides]
            if not slide_numbers:
                logger.warning("No valid slide numbers provided. Processing all slides.")
                slide_numbers = None
        
        result = {
            'file_path': pptx_path,
            'slides': [],
            'total_slides': total_slides,
            'slides_processed': 0,
            'tables_found': 0,
            'errors': []
        }
        
        slides_to_process = slide_numbers if slide_numbers else list(range(1, total_slides + 1))
        
        for slide_num in slides_to_process:
            try:
                slide_idx = slide_num - 1  # Convert to 0-indexed
                slide = prs.slides[slide_idx]
                
                logger.info(f"Processing slide {slide_num}/{total_slides}")
                
                slide_data = self._extract_slide_content(slide, slide_num, use_ocr, use_llm_formatting)
                result['slides'].append(slide_data)
                result['slides_processed'] += 1
                result['tables_found'] += len(slide_data.get('tables', []))
                
            except Exception as e:
                logger.error(f"Error processing slide {slide_num}: {e}")
                result['errors'].append(f"Slide {slide_num}: {str(e)}")
        
        return result
    
    def _extract_slide_content(self, slide, slide_num: int, use_ocr: bool, 
                               use_llm_formatting: bool) -> Dict[str, Any]:
        """Extract content from a single slide."""
        slide_data = {
            'slide_number': slide_num,
            'text': '',
            'tables': [],
            'images': [],
            'notes': '',
            'extraction_method': 'native',
            'ocr_used': False,
            'llm_formatted': False
        }
        
        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data['notes'] = slide.notes_slide.notes_text_frame.text
        
        # Extract text from shapes
        text_parts = []
        tables = []
        images = []
        
        for shape in slide.shapes:
            try:
                # Handle different shape types
                if shape.has_text_frame:
                    # Regular text boxes
                    text = self._extract_text_from_text_frame(shape.text_frame)
                    if text:
                        text_parts.append(text)
                
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    # Extract tables
                    table_data = self._extract_table_from_shape(shape)
                    if table_data:
                        tables.append(table_data)
                
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Handle images
                    if use_ocr and OCR_AVAILABLE:
                        image_data = self._extract_image_from_shape(shape)
                        if image_data:
                            images.append(image_data)
                            
                            # Perform OCR on image
                            try:
                                ocr_text = self._perform_ocr_on_image(image_data)
                                if ocr_text:
                                    text_parts.append(f"[Image OCR]: {ocr_text}")
                                    slide_data['ocr_used'] = True
                            except Exception as e:
                                logger.warning(f"OCR failed on image: {e}")
                
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Handle grouped shapes
                    grouped_text = self._extract_text_from_group(shape)
                    if grouped_text:
                        text_parts.append(grouped_text)
                
            except Exception as e:
                logger.warning(f"Error extracting from shape: {e}")
        
        # Combine text parts with proper formatting
        slide_data['text'] = self._format_text_parts(text_parts)
        slide_data['tables'] = tables
        slide_data['images'] = images
        
        # Use LLM for formatting if enabled
        if use_llm_formatting and OPENAI_AVAILABLE and self.openai_client:
            try:
                logger.info(f"Applying LLM formatting to slide {slide_num}")
                formatted_text = self._format_text_with_llm(slide_data['text'])
                slide_data['text'] = formatted_text
                slide_data['llm_formatted'] = True
            except Exception as e:
                logger.warning(f"LLM formatting failed: {e}")
        
        return slide_data
    
    def _extract_text_from_text_frame(self, text_frame) -> str:
        """Extract text from a text frame."""
        paragraphs = []
        for paragraph in text_frame.paragraphs:
            para_text = paragraph.text.strip()
            if para_text:
                paragraphs.append(para_text)
        return '\n'.join(paragraphs)
    
    def _extract_text_from_group(self, group_shape) -> str:
        """Extract text from grouped shapes."""
        text_parts = []
        for shape in group_shape.shapes:
            if shape.has_text_frame:
                text = self._extract_text_from_text_frame(shape.text_frame)
                if text:
                    text_parts.append(text)
        return '\n'.join(text_parts)
    
    def _extract_table_from_shape(self, shape) -> Optional[List[List[str]]]:
        """Extract table data from a table shape."""
        if not PANDAS_AVAILABLE:
            return None
        
        try:
            table = shape.table
            table_data = []
            
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_data.append(cell_text)
                table_data.append(row_data)
            
            return table_data
        except Exception as e:
            logger.error(f"Error extracting table: {e}")
            return None
    
    def _extract_image_from_shape(self, shape) -> Optional[bytes]:
        """Extract image data from a picture shape."""
        try:
            image = shape.image
            return image.blob
        except Exception as e:
            logger.error(f"Error extracting image: {e}")
            return None
    
    def _perform_ocr_on_image(self, image_data: bytes) -> str:
        """Perform OCR on image data."""
        if not OCR_AVAILABLE:
            return ""
        
        try:
            pil_image = Image.open(BytesIO(image_data))
            ocr_text = pytesseract.image_to_string(pil_image, config='--psm 6')
            return ocr_text
        except Exception as e:
            logger.error(f"OCR error: {e}")
            return ""
    
    def _format_text_parts(self, text_parts: List[str]) -> str:
        """Format text parts into readable content."""
        if not text_parts:
            return ""
        
        # Separate different sections
        formatted = []
        for i, part in enumerate(text_parts):
            if part.strip():
                formatted.append(part)
                # Add spacing between sections
                if i < len(text_parts) - 1:
                    formatted.append("")
        
        return '\n'.join(formatted)
    
    def _format_text_with_llm(self, text: str) -> str:
        """Use LLM to format text for better readability."""
        if not OPENAI_AVAILABLE or not self.openai_client:
            return text
        
        try:
            prompt = f"""You are a text formatting assistant. The following text was extracted from a PowerPoint slide where text elements can be positioned anywhere on the slide. Please format this text into a clean, readable structure with proper paragraphs, headings, and organization.

Extracted text:
{text}

Please format this text to be easy to read, with:
- Clear paragraphs
- Proper spacing
- Headings/sections where appropriate
- Bullet points or lists preserved
- Logical flow and organization

Formatted text:"""
            
            response = self.openai_client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that formats messy text into clean, readable content."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=2000,
                temperature=0.3
            )
            
            formatted_text = response.choices[0].message.content.strip()
            return formatted_text
            
        except Exception as e:
            logger.error(f"LLM formatting error: {e}")
            return text
    
    def save_text_content(self, result: Dict[str, Any], output_path: str):
        """Save extracted text content to a file."""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"PowerPoint Text Extraction Results\n")
            f.write(f"File: {result['file_path']}\n")
            f.write(f"Total Slides: {result['total_slides']}\n")
            f.write(f"Slides Processed: {result['slides_processed']}\n")
            f.write(f"Tables Found: {result['tables_found']}\n")
            f.write("=" * 50 + "\n\n")
            
            for slide in result['slides']:
                f.write(f"SLIDE {slide['slide_number']}\n")
                f.write("-" * 30 + "\n")
                
                if slide['text']:
                    f.write("TEXT CONTENT:\n")
                    f.write(slide['text'])
                    f.write("\n\n")
                
                if slide['notes']:
                    f.write("NOTES:\n")
                    f.write(slide['notes'])
                    f.write("\n\n")
                
                if slide['images']:
                    f.write(f"IMAGES FOUND: {len(slide['images'])}\n\n")
            
            if result['errors']:
                f.write("ERRORS:\n")
                for error in result['errors']:
                    f.write(f"- {error}\n")
        
        logger.info(f"Text content saved to: {output_path}")
    
    def save_tables_to_excel(self, result: Dict[str, Any], output_path: str):
        """Save extracted tables to Excel file."""
        if not PANDAS_AVAILABLE:
            logger.error("pandas not available. Cannot export tables to Excel.")
            return
        
        try:
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                table_count = 0
                
                for slide in result['slides']:
                    for i, table in enumerate(slide.get('tables', []), 1):
                        # Convert to DataFrame
                        df = pd.DataFrame(table)
                        
                        # Create sheet name (Excel has restrictions)
                        sheet_name = f"Slide{slide['slide_number']}_Table{i}"
                        # Remove invalid characters for Excel sheet names
                        invalid_chars = ['\\', '/', '?', '*', '[', ']', ':', "'"]
                        for char in invalid_chars:
                            sheet_name = sheet_name.replace(char, '_')
                        # Excel sheet name limit is 31 characters
                        if len(sheet_name) > 31:
                            sheet_name = sheet_name[:31]
                        
                        # Write to Excel
                        df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
                        table_count += 1
                
                if table_count == 0:
                    logger.warning("No tables found to export.")
                else:
                    logger.info(f"Exported {table_count} tables to: {output_path}")
                    
        except Exception as e:
            logger.error(f"Error saving tables to Excel: {e}")


def prompt_slide_selection(total_slides: int) -> Optional[List[int]]:
    """Prompt user for slide selection."""
    print(f"\nPresentation has {total_slides} slides.")
    print("Options:")
    print("  1. Extract all slides")
    print("  2. Extract specific slides")
    
    while True:
        choice = input("\nEnter your choice (1 or 2): ").strip()
        
        if choice == '1':
            return None  # All slides
        
        elif choice == '2':
            print("\nEnter slide numbers (separated by commas, e.g., 1,3,5):")
            slide_input = input("Slide numbers: ").strip()
            
            try:
                slide_numbers = [int(s.strip()) for s in slide_input.split(',')]
                # Validate slide numbers
                valid_numbers = [sn for sn in slide_numbers if 1 <= sn <= total_slides]
                
                if not valid_numbers:
                    print(f"Invalid slide numbers. Please enter numbers between 1 and {total_slides}.")
                    continue
                
                return valid_numbers
                
            except ValueError:
                print("Invalid input. Please enter numbers separated by commas.")
                continue
        
        else:
            print("Invalid choice. Please enter 1 or 2.")


def main():
    """Main function for command-line usage."""
    parser = argparse.ArgumentParser(description='Extract text from PowerPoint files with OCR and table extraction')
    parser.add_argument('pptx_path', help='Path to the PowerPoint file')
    parser.add_argument('-o', '--output-text', help='Output text file path')
    parser.add_argument('-e', '--output-excel', help='Output Excel file path for tables')
    parser.add_argument('--slides', nargs='+', type=int, help='Specific slide numbers to extract (1-indexed)')
    parser.add_argument('--no-ocr', action='store_true', help='Disable OCR for images')
    parser.add_argument('--no-llm', action='store_true', help='Disable LLM text formatting')
    parser.add_argument('--tesseract-path', help='Path to tesseract executable')
    parser.add_argument('--openai-key', help='OpenAI API key')
    parser.add_argument('--interactive', action='store_true', help='Interactive mode for slide selection')
    parser.add_argument('--verbose', '-v', action='store_true', help='Enable verbose logging')
    
    args = parser.parse_args()
    
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    # Initialize extractor
    extractor = PowerPointExtractor(
        tesseract_path=args.tesseract_path,
        openai_api_key=args.openai_key
    )
    
    try:
        # Determine slide selection
        slide_numbers = args.slides
        if args.interactive and not slide_numbers:
            # Load presentation to get total slides
            prs = Presentation(args.pptx_path)
            total_slides = len(prs.slides)
            slide_numbers = prompt_slide_selection(total_slides)
        
        # Extract content
        result = extractor.extract_text_from_pptx(
            pptx_path=args.pptx_path,
            slide_numbers=slide_numbers,
            use_ocr=not args.no_ocr,
            use_llm_formatting=not args.no_llm
        )
        
        # Determine output paths
        base_name = Path(args.pptx_path).stem
        
        if args.output_text:
            text_output = args.output_text
        else:
            text_output = f"{base_name}_extracted.txt"
        
        if args.output_excel:
            excel_output = args.output_excel
        else:
            excel_output = f"{base_name}_tables.xlsx"
        
        # Save results
        extractor.save_text_content(result, text_output)
        
        if result['tables_found'] > 0:
            extractor.save_tables_to_excel(result, excel_output)
        
        # Print summary
        print(f"\nExtraction completed!")
        print(f"Slides processed: {result['slides_processed']}/{result['total_slides']}")
        print(f"Tables found: {result['tables_found']}")
        print(f"Text output saved to: {text_output}")
        if result['tables_found'] > 0:
            print(f"Tables exported to: {excel_output}")
        
        if result['errors']:
            print(f"Errors encountered: {len(result['errors'])}")
            for error in result['errors']:
                print(f"  - {error}")
    
    except Exception as e:
        logger.error(f"Failed to process PowerPoint: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()

